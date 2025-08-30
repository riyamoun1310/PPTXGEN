# Backend - FastAPI setup

from fastapi import FastAPI, File, UploadFile, Form
from typing import Optional
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
import tempfile
import os
from pptx import Presentation
from pptx.util import Pt
import httpx
import json

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def apply_font_from_template(text_frame, template_placeholder):
    """
    Helper function to copy font style from a template placeholder
    to a new text frame.
    """
    # Get the font from the first paragraph of the template's placeholder
    template_font = template_placeholder.text_frame.paragraphs[0].font
    
    # Apply the template font properties to all paragraphs in the new text_frame
    for paragraph in text_frame.paragraphs:
        # Set font properties for each run in the paragraph
        for run in paragraph.runs:
            run.font.name = template_font.name
            run.font.size = template_font.size
            run.font.bold = template_font.bold
            run.font.italic = template_font.italic
            if template_font.color.rgb:
                run.font.color.rgb = template_font.color.rgb


async def get_slides_from_tavily(text: str, guidance: Optional[str], api_key: str):
    """
    Calls the Tavily API to split the input text into slides (title/content).
    Returns a list of dicts: [{"title": "...", "content": "..."}]
    """
    prompt = (
        "Split the following text into a list of slides for a PowerPoint presentation. "
        "Each slide should have a title and content. "
        f"{'Guidance: ' + guidance if guidance else ''}\n"
        f"Text:\n{text}\n"
        "Return as JSON: [{\"title\": \"...\", \"content\": \"...\"}]"
    )
    tavily_url = "https://api.tavily.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    data = {
        "model": "tavily-1",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": 1024
    }
    async with httpx.AsyncClient() as client:
        response = await client.post(tavily_url, headers=headers, json=data, timeout=60)
        if response.status_code != 200:
            raise RuntimeError(f"Tavily API error: {response.text}")
        result = response.json()
        content = result["choices"][0]["message"]["content"]
        slides = json.loads(content)
        if not isinstance(slides, list) or not all("title" in s and "content" in s for s in slides):
            raise ValueError("Tavily response format invalid.")
        return slides

@app.post("/generate")
async def generate_pptx(
    text: str = Form(...),
    guidance: str = Form(None),
    api_key: str = Form(...),
    pptx_file: Optional[UploadFile] = File(None)
):
    # --- LLM LOGIC ---
    try:
        slides_content = await get_slides_from_tavily(text, guidance, api_key)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=400)


    if pptx_file is not None:
        # --- TEMPLATE LOGIC ---
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(await pptx_file.read())
                tmp_path = tmp.name

            prs = Presentation(tmp_path)
            slide_layout = prs.slide_layouts[1]

            for item in slides_content:
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = item["title"]
                    template_title_placeholder = slide_layout.placeholders[0]
                    apply_font_from_template(slide.shapes.title.text_frame, template_title_placeholder)
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = item["content"]
                    template_content_placeholder = slide_layout.placeholders[1]
                    apply_font_from_template(content_placeholder.text_frame, template_content_placeholder)

            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as out_tmp:
                prs.save(out_tmp.name)
                out_path = out_tmp.name

            os.unlink(tmp_path)
            return FileResponse(out_path, filename="generated.pptx", media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        except Exception as e:
            return JSONResponse({"error": f"Template processing error: {str(e)}"}, status_code=500)
    else:
        # --- NO TEMPLATE LOGIC ---
        try:
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]
            for item in slides_content:
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]
                title.text = item["title"]
                body.text = item["content"]

            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                prs.save(tmp.name)
                tmp_path = tmp.name
            return FileResponse(tmp_path, filename="generated.pptx", media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        except Exception as e:
            return JSONResponse({"error": f"Presentation generation error: {str(e)}"}, status_code=500)


@app.get("/")
def root():
    return {"status": "Backend running"}
